"""Конвертация файлов ядра в сырой текст (PDF/DOCX/DOCM/DOC/PPTX).

Вход: `discover.DiscoveredFile`. Выход: `RawDocument` (сырой текст + при
наличии — постраничная разбивка для PDF и метаданные документа) для
`normalize.py`/`metadata.py`. Отказ конвертации → `ConversionError`
(источник кода INGEST-001, реестр `docs/dev/ERRORS.md`).

Зависимости: PyMuPDF (fitz) для PDF; ручной разбор OOXML (zip + XML) для
DOCX/DOCM — обходит строгую проверку content-type в python-docx, которая
падает на .docm (macro-enabled); `soffice --headless` для легаси .doc
(конвертация в .docx во временный каталог, затем тот же разбор OOXML);
python-pptx для PPTX.
"""
from __future__ import annotations

import subprocess
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from xml.etree import ElementTree as ET

import fitz  # PyMuPDF
from pptx import Presentation

from ariadna.ingest.config import SOFFICE_TIMEOUT_SEC

_W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
_CP_NS = "{http://schemas.openxmlformats.org/package/2006/metadata/core-properties}"
_DC_NS = "{http://purl.org/dc/elements/1.1/}"
_DCTERMS_NS = "{http://purl.org/dc/terms/}"


class ConversionError(Exception):
    """Не удалось получить текст из файла — источник кода INGEST-001."""


@dataclass
class RawDocument:
    """Сырой результат конвертации до нормализации/чанкинга."""

    text: str
    pages: list[str] | None  # заполнено только для PDF — нужно для чистки колонтитулов
    meta_title: str = ""
    meta_authors: list[str] | None = None
    meta_year: int | None = None


# ─── convert_pdf ─────────────────────────────────────────────────────────
# Назначение: извлекает текст постранично из PDF через PyMuPDF, забирает
#   метаданные документа (title/author/creationDate).
# Входные связи: путь к .pdf
# Выходные данные: RawDocument c pages=list[str] (по одной строке на страницу)
# Уровень: ✅ реализовано (A-02, worklogs/ingest.md#2026-07-03)
def convert_pdf(path: Path) -> RawDocument:
    try:
        doc = fitz.open(path)
    except Exception as exc:  # битый/зашифрованный PDF
        raise ConversionError(f"не удалось открыть PDF: {exc}") from exc
    if doc.is_encrypted:
        raise ConversionError("PDF зашифрован")
    pages = [page.get_text() for page in doc]
    meta = doc.metadata or {}
    authors = _split_authors(meta.get("author", ""))
    year = _year_from_pdf_date(meta.get("creationDate", ""))
    return RawDocument(
        text="\n".join(pages),
        pages=pages,
        meta_title=(meta.get("title") or "").strip(),
        meta_authors=authors,
        meta_year=year,
    )


# ─── convert_ooxml_word ──────────────────────────────────────────────────
# Назначение: извлекает текст и core-свойства из DOCX/DOCM без python-docx
#   (обходит его строгую валидацию content-type, ломающуюся на .docm).
# Входные связи: путь к .docx/.docm
# Выходные данные: RawDocument (pages=None — постраничности в OOXML нет)
# Уровень: ✅ реализовано (A-02, worklogs/ingest.md#2026-07-03)
def convert_ooxml_word(path: Path) -> RawDocument:
    try:
        with zipfile.ZipFile(path) as zf:
            text = _extract_ooxml_paragraphs(zf)
            title, authors, year = _extract_ooxml_core_props(zf)
    except (zipfile.BadZipFile, KeyError) as exc:
        raise ConversionError(f"не удалось разобрать OOXML: {exc}") from exc
    return RawDocument(text=text, pages=None, meta_title=title, meta_authors=authors, meta_year=year)


# ─── convert_legacy_doc ──────────────────────────────────────────────────
# Назначение: конвертирует старый бинарный .doc в .docx через headless
#   libreoffice, затем разбирает результат тем же OOXML-парсером.
# Входные связи: путь к .doc; системный soffice (проверено — установлен)
# Выходные данные: RawDocument, как и convert_ooxml_word
# Уровень: ✅ реализовано (A-02, worklogs/ingest.md#2026-07-03)
def convert_legacy_doc(path: Path) -> RawDocument:
    with tempfile.TemporaryDirectory() as out_dir, tempfile.TemporaryDirectory() as profile_dir:
        cmd = [
            "soffice",
            "--headless",
            "--convert-to",
            "docx",
            "--outdir",
            out_dir,
            f"-env:UserInstallation=file://{profile_dir}",
            str(path),
        ]
        try:
            subprocess.run(
                cmd, check=True, timeout=SOFFICE_TIMEOUT_SEC,
                capture_output=True,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as exc:
            raise ConversionError(f"libreoffice не сконвертировал .doc: {exc}") from exc
        converted = Path(out_dir) / f"{path.stem}.docx"
        if not converted.exists():
            raise ConversionError("libreoffice не создал .docx-результат")
        return convert_ooxml_word(converted)


# ─── convert_pptx ─────────────────────────────────────────────────────────
# Назначение: извлекает текст слайдов (все текстовые фреймы) и core-свойства
#   презентации через python-pptx.
# Входные связи: путь к .pptx
# Выходные данные: RawDocument (pages=None — режем по чанкингу, не по слайдам)
# Уровень: ✅ реализовано (A-02, worklogs/ingest.md#2026-07-03)
def convert_pptx(path: Path) -> RawDocument:
    try:
        pres = Presentation(str(path))
    except Exception as exc:
        raise ConversionError(f"не удалось открыть PPTX: {exc}") from exc
    chunks: list[str] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs)
                if line.strip():
                    chunks.append(line)
    props = pres.core_properties
    authors = _split_authors(props.author or "")
    year = props.created.year if props.created else None
    return RawDocument(
        text="\n".join(chunks),
        pages=None,
        meta_title=(props.title or "").strip(),
        meta_authors=authors,
        meta_year=year,
    )


# Назначение: текст всех непустых параграфов word/document.xml по порядку.
# Уровень: ✅ реализовано (A-02, worklogs/ingest.md#2026-07-03)
def _extract_ooxml_paragraphs(zf: zipfile.ZipFile) -> str:
    with zf.open("word/document.xml") as f:
        root = ET.parse(f).getroot()
    paragraphs = []
    for p in root.iter(f"{_W_NS}p"):
        text = "".join(t.text or "" for t in p.iter(f"{_W_NS}t"))
        if text.strip():
            paragraphs.append(text)
    return "\n".join(paragraphs)


# Назначение: title/creator/created год из docProps/core.xml (если part есть).
# Уровень: ✅ реализовано (A-02, worklogs/ingest.md#2026-07-03)
def _extract_ooxml_core_props(zf: zipfile.ZipFile) -> tuple[str, list[str], int | None]:
    if "docProps/core.xml" not in zf.namelist():
        return "", [], None
    with zf.open("docProps/core.xml") as f:
        root = ET.parse(f).getroot()
    title = (root.findtext(f"{_DC_NS}title") or "").strip()
    creator = (root.findtext(f"{_DC_NS}creator") or "").strip()
    created_raw = root.findtext(f"{_DCTERMS_NS}created") or ""
    year = None
    if len(created_raw) >= 4 and created_raw[:4].isdigit():
        year = int(created_raw[:4])
    return title, _split_authors(creator), year


# Назначение: один текст свойств документа → список имён (разделители ; , / \n).
# Уровень: ✅ реализовано (A-02, worklogs/ingest.md#2026-07-03)
def _split_authors(raw: str) -> list[str]:
    for sep in (";", "/", "\n"):
        raw = raw.replace(sep, ",")
    return [name.strip() for name in raw.split(",") if name.strip()]


# Назначение: PDF creationDate вида «D:20110316161454+03'00'» → год или None.
# Уровень: ✅ реализовано (A-02, worklogs/ingest.md#2026-07-03)
def _year_from_pdf_date(raw: str) -> int | None:
    digits = raw.removeprefix("D:")
    if len(digits) >= 4 and digits[:4].isdigit():
        return int(digits[:4])
    return None
