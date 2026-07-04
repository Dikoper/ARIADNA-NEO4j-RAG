"""build_pptx.py — генератор презентации для жюри хакатона (модуль docs, задача A-19b).

Назначение: собрать docs/submission/presentation/ariadna.pptx из данных SLIDES ниже.
Весь текст слайдов вынесен в структуры данных в начале файла — чтобы PM мог попросить
поправить формулировки без разбора рендер-кода. Рендер-функции (ниже раздела DATA)
просто раскладывают эти данные по слайдам в сдержанном деловом стиле:
тёмно-синий/белый, без пёстрых цветов, крупные цифры.

Зависимость: python-pptx (уже в pyproject.toml [project].dependencies).

Запуск:
    .venv/bin/python docs/submission/presentation/build_pptx.py

PDF (опционально, если есть LibreOffice):
    soffice --headless --convert-to pdf --outdir docs/submission/presentation \
        docs/submission/presentation/ariadna.pptx

Это скрипт модуля docs — не часть src/ariadna, пре-комментарии CONVENTIONS.md
(обязательны для src/) на него не распространяются; докстринг функций оставлен
для читаемости и повторного использования PM/оркестратором.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUTPUT_PATH = Path(__file__).parent / "ariadna.pptx"

# ─────────────────────────────────────────────────────────────────────────────
# DATA — содержимое слайдов. Правки текста презентации делаются ТОЛЬКО здесь.
# ─────────────────────────────────────────────────────────────────────────────

FOOTER_BRAND = "Ариадна · хакатон Норникеля «Научный клубок»"

SLIDES: list[dict] = [
    {
        "type": "title",
        "title": "Ариадна",
        "subtitle": "Карта знаний R&D горно-металлургической отрасли",
        "meta": [
            "Хакатон Норникеля — задача «Научный клубок»",
            "Команда «Ариадна» · 4 июля 2026",
        ],
    },
    {
        "type": "bullets",
        "title": "Проблема",
        "bullets": [
            "1 453 документа, 4,9 ГБ — статьи, обзоры, отчёты, доклады на русском "
            "и английском языках",
            "Знания рассеяны по корпусу: опыт, числовые параметры, противоречия "
            "ищутся вручную — ни один эксперт не может перечитать всё",
            "Непонятно, где реально нет исследований: «пробелы» не видны, пока "
            "на них случайно не наткнёшься",
            "Вопросы звучат естественно («какие методы обессоливания подходят, "
            "если …») — а корпус не структурирован для такого поиска",
        ],
    },
    {
        "type": "flow",
        "title": "Решение одним взглядом",
        "boxes": [
            "Документы\n(1 453 файла)",
            "Извлечение\nLLM по онтологии\n+ правила для чисел",
            "Граф знаний\nNeo4j",
            "Вопрос на\nестественном языке",
            "Ответ\nс цитатами + подграф\n+ карта пробелов",
        ],
        "caption": "Каждый факт в графе — с провенансом до чанка-источника",
    },
    {
        "type": "bullets",
        "title": "Архитектура: двойной граф",
        "bullets": [
            "Лексический граф: Document → Chunk + векторный индекс "
            "(Qwen3-Embedding) — страховочный RAG",
            "Сущностный граф: онтология задания — 8 типов сущностей, 6 типов "
            "связей, два хаба — Experiment и TechSolution/Process",
            "Гибридный retrieval: граф + вектор, единое ранжирование чанков",
            "Роутер запросов без LLM: намерение определяют правила → "
            "шаблонные Cypher-запросы (свободный text2cypher запрещён)",
        ],
    },
    {
        "type": "stats",
        "title": "Масштаб в цифрах",
        "stats": [
            ("177", "документов ядра"),
            ("9 580", "чанков"),
            ("23,6 тыс.", "сущностей"),
            ("25 063", "связей графа"),
            ("8 179", "числовых ограничений"),
            ("175 / 177", "документов с гео-разметкой"),
            ("620", "автотестов (+3 xfail)"),
        ],
    },
    {
        "type": "bullets",
        "title": "Точность чисел — не дело случая",
        "kicker": "ИНВАРИАНТ СИСТЕМЫ",
        "bullets": [
            "Концентрации, температуры, расходы и другие числовые параметры "
            "извлекает ТОЛЬКО регекс-нормализатор — никогда LLM",
            "LLM определяет, где искать; финальное число и единицу измерения "
            "проверяет детерминированное правило",
            "Ошибка в «сульфаты ≤300 мг/л» — не опечатка, а неверный "
            "технологический вывод для инженера",
            "8 179 числовых ограничений в графе — каждое привязано к "
            "чанку-источнику",
        ],
    },
    {
        "type": "bullets",
        "title": "4 запроса жюри — отвечают все",
        "bullets": [
            "Q1. Обессоливание воды (сульфаты/хлориды/Ca/Mg/Na 200–300 мг/л, "
            "сухой остаток ≤1000 мг/дм³) — применимые методы с граничными "
            "концентрациями",
            "Q2. Циркуляция католита при электроэкстракции никеля — решения "
            "мировой практики и оптимальная скорость потока",
            "Q3. Распределение Au, Ag и МПГ между штейном и шлаком за 5 лет — "
            "эксперименты и публикации по теме",
            "Q4. Закачка шахтных вод в глубокие горизонты — честный пробел "
            "в корпусе + смежные альтернативы",
        ],
        "note": "Каждый ответ — с цитатами до конкретного чанка-источника",
    },
    {
        "type": "bullets",
        "title": "Карта пробелов",
        "kicker": "ГЛАВНАЯ ФИЧА",
        "bullets": [
            "Gap-матрица «материал × процесс»: показывает сочетания, для "
            "которых в корпусе нет экспериментов",
            "Q4 (закачка шахтных вод) — теме физически нечего ответить; "
            "система честно говорит «в корпусе не найдено» и предлагает "
            "смежные альтернативы",
            "Гео-разметка вскрывает перекос практики: только в РФ (only_ru) — "
            "44 темы, только за рубежом (only_foreign) — 225 тем",
            "Пробел — не баг, а сигнал: куда направить R&D или литобзор",
        ],
    },
    {
        "type": "bullets",
        "title": "Верификация и работа эксперта",
        "bullets": [
            "Каждый факт несёт provenance: source, updated_at, confidence",
            "Противоречия (contradicts) подсвечиваются в подграфе ответа, "
            "а не прячутся",
            "Эксперт правит граф напрямую в Neo4j Browser — без "
            "переизвлечения корпуса",
            "Версионирование-минимум готово к аудиту решений",
        ],
    },
    {
        "type": "bullets",
        "title": "Стек — полностью локальный",
        "bullets": [
            "Neo4j — граф знаний + векторный индекс",
            "Qwen3.5-35B-A3B через Ollama (int4) — извлечение и синтез "
            "ответов, локально на DGX Spark",
            "Qwen3-Embedding-0.6B — эмбеддинги RU/EN",
            "Ни документ, ни вопрос, ни ответ не покидают контур — внешних "
            "API нет",
        ],
    },
    {
        "type": "bullets",
        "title": "Что дальше",
        "bullets": [
            "Рекомендации: похожие кейсы, эксперты по теме, смежные направления",
            "Литобзор с автоматическим консенсусом/разногласиями источников",
            "JSON-LD фасад для интеграции с внешними системами",
            "Масштабирование извлечения на весь корпус (1 453 документа)",
        ],
    },
    {
        "type": "closing",
        "title": "Демо",
        "bullets": [
            "Streamlit-чат: вопрос → ответ с цитатами → подграф → карта пробелов",
            "Видео и материалы: disk.yandex.ru/d/Vo0tJWsdzO_sIg",
            "Код: github.com/Dikoper/ARIADNA-NEO4j-RAG",
        ],
        "thanks": "Спасибо!",
    },
]

# ─────────────────────────────────────────────────────────────────────────────
# STYLE — сдержанная деловая палитра. Без пёстрых цветов.
# ─────────────────────────────────────────────────────────────────────────────

NAVY_DARK = RGBColor(0x0B, 0x1B, 0x2E)  # фон титульного/финального слайда
NAVY = RGBColor(0x10, 0x24, 0x3E)  # основной текст заголовков
ACCENT = RGBColor(0x3E, 0x66, 0x92)  # приглушённый акцент (цифры, буллеты)
GRAY = RGBColor(0x51, 0x5E, 0x70)  # вторичный текст
GRAY_LIGHT = RGBColor(0x8A, 0x95, 0xA5)  # футер, метки
LINE = RGBColor(0xD8, 0xDE, 0xE6)  # тонкие разделители
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF6, 0xF8, 0xFA)  # плашки на светлом фоне

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


def _set_text(
    tf,
    text: str,
    *,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font: str = FONT,
    line_spacing: float | None = None,
) -> None:
    """Назначение: заполнить первый параграф текстового фрейма с заданным стилем.
    Уровень: реализовано (A-19b)."""
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run() if not p.runs else p.runs[0]
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _add_footer(slide, page_no: int, *, dark_bg: bool = False) -> None:
    """Назначение: тонкая нижняя строка — бренд + номер слайда.
    Уровень: реализовано (A-19b)."""
    color = GRAY_LIGHT if not dark_bg else RGBColor(0x8F, 0xA0, 0xBA)
    box = _add_textbox(
        None,
        Emu(0),
        SLIDE_H - Inches(0.42),
        SLIDE_W - MARGIN,
        Inches(0.32),
        slide=slide,
    )
    tf = box.text_frame
    tf.margin_left = MARGIN
    tf.margin_right = MARGIN
    _set_text(tf, f"{FOOTER_BRAND}    ·    {page_no:02d}", size=10, color=color)


def _add_textbox(prs, left, top, width, height, *, slide):
    """Назначение: добавить пустой текстовый фрейм на переданный слайд.
    Уровень: реализовано (A-19b)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    return box


def _blank_slide(prs):
    """Назначение: пустой слайд (layout без плейсхолдеров) — полный контроль над версткой.
    Уровень: реализовано (A-19b)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_bg(slide, color: RGBColor) -> None:
    """Назначение: залить фон слайда сплошным цветом.
    Уровень: реализовано (A-19b)."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _add_title(slide, text: str, *, color: RGBColor = NAVY, kicker: str | None = None) -> None:
    """Назначение: заголовок слайда + тонкая акцентная линия под ним + опц. метка-киккер.
    Уровень: реализовано (A-19b)."""
    if kicker:
        kbox = slide.shapes.add_textbox(MARGIN, Inches(0.42), Inches(10), Inches(0.35))
        _set_text(kbox.text_frame, kicker, size=13, color=ACCENT, bold=True)
    top = Inches(0.75) if kicker else Inches(0.55)
    tbox = slide.shapes.add_textbox(MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.9))
    _set_text(tbox.text_frame, text, size=32, color=color, bold=True)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, top + Inches(0.85), Inches(1.1), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def render_title_slide(prs, data: dict, page_no: int) -> None:
    """Назначение: титульный слайд — тёмно-синий фон, крупное имя проекта.
    Уровень: реализовано (A-19b)."""
    slide = _blank_slide(prs)
    _fill_bg(slide, NAVY_DARK)

    title_box = slide.shapes.add_textbox(
        MARGIN, Inches(2.55), SLIDE_W - 2 * MARGIN, Inches(1.3)
    )
    _set_text(title_box.text_frame, data["title"], size=60, color=WHITE, bold=True)

    sub_box = slide.shapes.add_textbox(
        MARGIN, Inches(3.65), SLIDE_W - 2 * MARGIN, Inches(0.8)
    )
    _set_text(sub_box.text_frame, data["subtitle"], size=24, color=RGBColor(0xC9, 0xD6, 0xE8))

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(2.35), Inches(1.4), Pt(4)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()

    meta_top = Inches(6.3)
    for i, line in enumerate(data["meta"]):
        mbox = slide.shapes.add_textbox(
            MARGIN, meta_top + Inches(0.36) * i, SLIDE_W - 2 * MARGIN, Inches(0.35)
        )
        _set_text(mbox.text_frame, line, size=14, color=RGBColor(0x9A, 0xAC, 0xC7))


def render_bullets_slide(prs, data: dict, page_no: int) -> None:
    """Назначение: стандартный слайд «заголовок + список тезисов» (+ опц. заметка внизу).
    Уровень: реализовано (A-19b)."""
    slide = _blank_slide(prs)
    _fill_bg(slide, WHITE)
    _add_title(slide, data["title"], kicker=data.get("kicker"))

    top = Inches(2.15) if not data.get("kicker") else Inches(2.15)
    body = slide.shapes.add_textbox(MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(4.3))
    tf = body.text_frame
    tf.word_wrap = True
    bullets = data["bullets"]
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(18)
        p.line_spacing = 1.12
        run = p.add_run()
        run.text = f"●  {text}"
        run.font.size = Pt(19)
        run.font.color.rgb = NAVY
        run.font.name = FONT

    if data.get("note"):
        note_box = slide.shapes.add_textbox(
            MARGIN, SLIDE_H - Inches(1.15), SLIDE_W - 2 * MARGIN, Inches(0.5)
        )
        _set_text(note_box.text_frame, data["note"], size=14, color=GRAY, bold=False)

    _add_footer(slide, page_no)


def render_flow_slide(prs, data: dict, page_no: int) -> None:
    """Назначение: горизонтальная схема «документы → … → ответ» из прямоугольников и стрелок.
    Уровень: реализовано (A-19b)."""
    slide = _blank_slide(prs)
    _fill_bg(slide, WHITE)
    _add_title(slide, data["title"])

    boxes = data["boxes"]
    n = len(boxes)
    box_h = Inches(1.7)
    arrow_w = Inches(0.45)
    gap = Inches(0.12)
    total_margin = 2 * MARGIN
    available = SLIDE_W - total_margin - arrow_w * (n - 1) - gap * 2 * (n - 1)
    box_w = Emu(int(available / n))
    top = Inches(2.9)

    x = MARGIN
    for i, label in enumerate(boxes):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY if i in (0, n - 1) else BG_LIGHT
        shape.line.color.rgb = LINE
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_color = WHITE if i in (0, n - 1) else NAVY
        lines = label.split("\n")
        for j, line in enumerate(lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.size = Pt(15 if j == 0 else 12.5)
            run.font.bold = j == 0
            run.font.color.rgb = text_color
            run.font.name = FONT
        x += box_w + gap

        if i < n - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x,
                top + box_h / 2 - Inches(0.18),
                arrow_w,
                Inches(0.36),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()
            x += arrow_w + gap

    caption = slide.shapes.add_textbox(
        MARGIN, top + box_h + Inches(0.5), SLIDE_W - 2 * MARGIN, Inches(0.6)
    )
    _set_text(caption.text_frame, data["caption"], size=16, color=GRAY, align=PP_ALIGN.CENTER)

    _add_footer(slide, page_no)


def render_stats_slide(prs, data: dict, page_no: int) -> None:
    """Назначение: сетка плашек «крупная цифра + подпись» для слайда масштаба.
    Уровень: реализовано (A-19b)."""
    slide = _blank_slide(prs)
    _fill_bg(slide, WHITE)
    _add_title(slide, data["title"])

    stats = data["stats"]
    cols = 4
    rows = (len(stats) + cols - 1) // cols
    gap = Inches(0.25)
    top0 = Inches(2.25)
    tile_w = Emu(int((SLIDE_W - 2 * MARGIN - gap * (cols - 1)) / cols))
    tile_h = Inches(1.95)

    for idx, (value, label) in enumerate(stats):
        row, col = divmod(idx, cols)
        x = MARGIN + col * (tile_w + gap)
        y = top0 + row * (tile_h + gap)
        tile = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, tile_w, tile_h)
        tile.fill.solid()
        tile.fill.fore_color.rgb = BG_LIGHT
        tile.line.color.rgb = LINE
        tile.line.width = Pt(1)
        tf = tile.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = value
        r1.font.size = Pt(30)
        r1.font.bold = True
        r1.font.color.rgb = NAVY
        r1.font.name = FONT

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(13)
        r2.font.color.rgb = GRAY
        r2.font.name = FONT

    _add_footer(slide, page_no)


def render_closing_slide(prs, data: dict, page_no: int) -> None:
    """Назначение: финальный слайд — тёмно-синий фон, демо + благодарность.
    Уровень: реализовано (A-19b)."""
    slide = _blank_slide(prs)
    _fill_bg(slide, NAVY_DARK)

    title_box = slide.shapes.add_textbox(MARGIN, Inches(1.1), SLIDE_W - 2 * MARGIN, Inches(1.0))
    _set_text(title_box.text_frame, data["title"], size=48, color=WHITE, bold=True)

    body = slide.shapes.add_textbox(MARGIN, Inches(2.5), SLIDE_W - 2 * MARGIN, Inches(2.0))
    tf = body.text_frame
    tf.word_wrap = True
    for i, text in enumerate(data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        run = p.add_run()
        run.text = f"●  {text}"
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0xC9, 0xD6, 0xE8)
        run.font.name = FONT

    thanks_box = slide.shapes.add_textbox(MARGIN, Inches(5.6), SLIDE_W - 2 * MARGIN, Inches(1.0))
    _set_text(thanks_box.text_frame, data["thanks"], size=36, color=WHITE, bold=True)

    _add_footer(slide, page_no, dark_bg=True)


RENDERERS = {
    "title": render_title_slide,
    "bullets": render_bullets_slide,
    "flow": render_flow_slide,
    "stats": render_stats_slide,
    "closing": render_closing_slide,
}


def build(output_path: Path = OUTPUT_PATH) -> Path:
    """Назначение: собрать презентацию из SLIDES и сохранить .pptx на диск.
    Входные связи: SLIDES (данные слайдов, см. раздел DATA выше)
    Выходные данные: путь к сохранённому .pptx
    Уровень: реализовано (A-19b)."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for page_no, slide_data in enumerate(SLIDES, start=1):
        renderer = RENDERERS[slide_data["type"]]
        renderer(prs, slide_data, page_no)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = build()
    print(f"Презентация сохранена: {path} ({len(SLIDES)} слайдов)")
