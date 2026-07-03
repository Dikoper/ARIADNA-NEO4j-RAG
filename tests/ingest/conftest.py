"""Фикстуры для тестов ingest: генерация мини-PDF/DOCX/PPTX с известным содержимым.

Файлы пишутся в tests/ingest/fixtures/ (создаются программно при сборе тестов,
без ручного бинарного контента). Настоящие данные data/ не читаются здесь
(смоук на реальных data/processed/*.jsonl — отдельный тестовый модуль).
"""
from __future__ import annotations

import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path

import fitz  # PyMuPDF
import pytest
from pptx import Presentation

HAS_SOFFICE = shutil.which("soffice") is not None or shutil.which("libreoffice") is not None

FIXTURES_DIR = Path(__file__).parent / "fixtures"

# Base14 "helv" (Helvetica) в PyMuPDF не поддерживает кириллицу (WinAnsi-подобная
# кодировка режет русские буквы в "·") — для PDF-фикстур с русским текстом нужен
# встроенный TTF-шрифт с кириллицей. DejaVu Sans — стандартный пакет на Ubuntu.
_CYRILLIC_FONT_PATH = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"


def _insert_cyrillic_text(page, point, text: str, fontsize: float = 8) -> None:
    """Вставляет текст с поддержкой кириллицы через встроенный DejaVu Sans.

    Использует insert_textbox (не insert_text): insert_text не переносит строки
    и молча обрезает содержимое по правому краю страницы — на длинном тексте
    без явных "\\n" это теряет большую часть содержимого (проверено эмпирически).
    insert_textbox, в свою очередь, если текст совсем не помещается в rect
    (сильно отрицательный остаток), тоже пишет 0 символов — поэтому margin
    маленький и fontsize уменьшен, чтобы длинные фикстуры (~6-8 KB) помещались
    полностью (тоже проверено эмпирически подбором).
    """
    page.insert_font(fontname="F0", fontfile=_CYRILLIC_FONT_PATH)
    margin = 36
    rect = fitz.Rect(margin, margin, page.rect.width - margin, page.rect.height - margin)
    rc = page.insert_textbox(rect, text, fontsize=fontsize, fontname="F0")
    assert rc >= 0, (
        f"insert_textbox: текст не поместился в фикстуру (rc={rc}); "
        "уменьшите объём текста или fontsize в conftest.py"
    )

# Текст со спецсимволами разных категорий: индексы, неравенство, температура,
# единицы, диапазон с en-dash, смешанный RU/EN — все требования паспорта ingest.
SPECIAL_CHARS_SENTENCE = (
    "Концентрация H2SO4 (серной кислоты) не должна превышать значений: "
    "содержание SO₂ и CO₂ контролируется методом titration. "
    "Предел ≤ 200 мг/л при температуре не выше 80°C. "
    "Диапазон концентрации сульфатов составляет 200–300 мг/дм³ согласно ГОСТ. "
    "This process uses electrowinning (электроэкстракция) at low pH values."
)


# ─── special_chars_sentence ─────────────────────────────────────────────
# Назначение: доступ к эталонной строке спецсимволов из тестов других модулей
#   (без импорта tests.ingest как пакета — избегаем зависимости от __init__.py).
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def special_chars_sentence() -> str:
    return SPECIAL_CHARS_SENTENCE


# ─── has_soffice ─────────────────────────────────────────────────────────
# Назначение: доступность soffice/libreoffice в текущем окружении — тесты
#   .doc-конвертации сами себе строят фикстуру этим же путём и пропускаются,
#   если бинарник недоступен (см. worklogs/ingest.md, открытый вопрос).
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def has_soffice() -> bool:
    return HAS_SOFFICE


# ─── _write_ooxml_docx ──────────────────────────────────────────────────
# Назначение: собирает минимальный валидный .docx (word/document.xml +
#   docProps/core.xml + служебные части) с заданными параграфами и мета.
# Уровень: ✅ реализовано (A-02 tests)
def _write_ooxml_docx(
    path: Path,
    paragraphs: list[str],
    title: str = "",
    creator: str = "",
    created: str = "",
) -> None:
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        '<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>'
        '</Types>'
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="word/document.xml"/>'
        '<Relationship Id="rId2" '
        'Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" '
        'Target="docProps/core.xml"/>'
        '</Relationships>'
    )
    body_paragraphs = "".join(
        f'<w:p><w:r><w:t xml:space="preserve">{p}</w:t></w:r></w:p>' for p in paragraphs
    )
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f'<w:body>{body_paragraphs}</w:body></w:document>'
    )
    core_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<cp:coreProperties '
        'xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" '
        'xmlns:dc="http://purl.org/dc/elements/1.1/" '
        'xmlns:dcterms="http://purl.org/dc/terms/">'
        f'<dc:title>{title}</dc:title><dc:creator>{creator}</dc:creator>'
        f'{f"<dcterms:created>{created}</dcterms:created>" if created else ""}'
        '</cp:coreProperties>'
    )
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("word/document.xml", document_xml)
        zf.writestr("docProps/core.xml", core_xml)


# ─── docx_known_content ──────────────────────────────────────────────────
# Назначение: .docx с известными параграфами (включая спецсимволы) + мета.
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def docx_known_content() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "known_content.docx"
    _write_ooxml_docx(
        path,
        paragraphs=["Введение в тему исследования.", SPECIAL_CHARS_SENTENCE, "Заключение работы 2019."],
        title="Отчёт о лабораторных испытаниях",
        creator="Иванов И.И.; Петров П.П.",
        created="2019-05-01T00:00:00Z",
    )
    return path


# ─── docm_macro_enabled ───────────────────────────────────────────────────
# Назначение: .docm (macro-enabled) — python-docx падает на нём из-за строгой
#   проверки content-type; наш OOXML-парсер должен работать одинаково.
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def docm_macro_enabled() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "macro_enabled.docm"
    _write_ooxml_docx(path, paragraphs=["Текст документа с макросами."], title="", creator="")
    return path


# ─── docx_corrupted ────────────────────────────────────────────────────────
# Назначение: файл с расширением .docx, не являющийся валидным zip/OOXML.
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def docx_corrupted() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "corrupted.docx"
    path.write_bytes(b"this is not a zip file at all")
    return path


# ─── docx_empty ────────────────────────────────────────────────────────────
# Назначение: валидный .docx без единого параграфа текста (пустой документ).
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def docx_empty() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "empty.docx"
    _write_ooxml_docx(path, paragraphs=[], title="", creator="")
    return path


# ─── doc_legacy_from_soffice ────────────────────────────────────────────────
# Назначение: настоящий бинарный .doc, полученный конвертацией docx_known_content
#   через тот же soffice --headless, что использует convert_legacy_doc в коде
#   (открытый вопрос из worklogs/ingest.md: доступность soffice на целевом
#   окружении — здесь только фиксируем, что тест сам собирает фикстуру этим
#   же путём, если soffice есть в текущем окружении).
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def doc_legacy_from_soffice(docx_known_content: Path) -> Path:
    if not HAS_SOFFICE:
        pytest.skip("soffice/libreoffice недоступен в этом окружении")
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    out_path = FIXTURES_DIR / "known_content.doc"
    with tempfile.TemporaryDirectory() as out_dir, tempfile.TemporaryDirectory() as profile_dir:
        cmd = [
            "soffice", "--headless", "--convert-to", "doc",
            "--outdir", out_dir,
            f"-env:UserInstallation=file://{profile_dir}",
            str(docx_known_content),
        ]
        subprocess.run(cmd, check=True, timeout=120, capture_output=True)
        converted = Path(out_dir) / f"{docx_known_content.stem}.doc"
        assert converted.exists(), "soffice не создал .doc — фикстуру собрать не удалось"
        shutil.copy(converted, out_path)
    return out_path


# ─── pptx_known_content ────────────────────────────────────────────────────
# Назначение: .pptx с двумя слайдами известного текста через python-pptx.
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def pptx_known_content() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "known_content.pptx"
    pres = Presentation()
    layout = pres.slide_layouts[6]  # пустой макет
    slide1 = pres.slides.add_slide(layout)
    box1 = slide1.shapes.add_textbox(0, 0, 5000000, 1000000)
    box1.text_frame.text = "Первый слайд доклада."
    slide2 = pres.slides.add_slide(layout)
    box2 = slide2.shapes.add_textbox(0, 0, 5000000, 1000000)
    box2.text_frame.text = SPECIAL_CHARS_SENTENCE
    pres.core_properties.title = "Доклад на конференции"
    pres.core_properties.author = "Сидоров С.С."
    pres.save(str(path))
    return path


# ─── pdf_special_chars ──────────────────────────────────────────────────────
# Назначение: многостраничный PDF с текстовым слоем, спецсимволами и объёмом,
#   достаточным для нескольких чанков (~3-4 KB текста).
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def pdf_special_chars() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "special_chars.pdf"
    doc = fitz.open()
    body = (SPECIAL_CHARS_SENTENCE + " ") * 20  # достаточно текста для >1 чанка
    for _ in range(2):
        page = doc.new_page()
        _insert_cyrillic_text(page, (72, 72), body)
    doc.set_metadata({"title": "Спецсимволы 2021", "author": "Кузнецов К.К."})
    doc.save(str(path))
    doc.close()
    return path


# ─── pdf_with_boilerplate ───────────────────────────────────────────────────
# Назначение: PDF с повторяющимся колонтитулом (заголовок) и футером
#   (номер страницы) на каждой из 6 страниц — должны быть вычищены.
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def pdf_with_boilerplate() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "with_boilerplate.pdf"
    doc = fitz.open()
    header = "Журнал металлургических исследований №3"
    for i in range(6):
        page = doc.new_page()
        body = f"Уникальный текст страницы номер {i}. " + SPECIAL_CHARS_SENTENCE
        text = f"{header}\n{body}\n{i + 1}"
        _insert_cyrillic_text(page, (72, 72), text)
    doc.save(str(path))
    doc.close()
    return path


# ─── pdf_no_text_layer ──────────────────────────────────────────────────────
# Назначение: PDF без текстового слоя (только векторная графика) — эмулирует
#   скан без OCR; должен уйти в skip с INGEST-002.
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def pdf_no_text_layer() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "no_text_layer.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.draw_rect(fitz.Rect(50, 50, 200, 200))  # только графика, ни одного символа текста
    doc.save(str(path))
    doc.close()
    return path


# ─── pdf_too_short ───────────────────────────────────────────────────────────
# Назначение: PDF с текстовым слоем короче порога MIN_TEXT_CHARS (200 симв.).
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def pdf_too_short() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "too_short.pdf"
    doc = fitz.open()
    page = doc.new_page()
    _insert_cyrillic_text(page, (72, 72), "Короткий текст.")
    doc.save(str(path))
    doc.close()
    return path


# ─── pdf_corrupted ───────────────────────────────────────────────────────────
# Назначение: файл с расширением .pdf, не являющийся валидным PDF.
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def pdf_corrupted() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "corrupted.pdf"
    path.write_bytes(b"%PDF-1.4 this is not really a pdf body, just garbage bytes")
    return path


# ─── unsupported_format_file ─────────────────────────────────────────────────
# Назначение: файл неподдерживаемого формата (.rar) — для INGEST-001 через discover.
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture(scope="session")
def unsupported_format_file() -> Path:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    path = FIXTURES_DIR / "archive.rar"
    path.write_bytes(b"Rar!\x1a\x07\x01\x00fake archive bytes")
    return path


# ─── core_fixture_dir ─────────────────────────────────────────────────────────
# Назначение: собирает мини-корпус data_dir/<CORE_FOLDER>/* для прогона
#   run_pipeline целиком (discover→convert→normalize→metadata→chunk→jsonl).
# Уровень: ✅ реализовано (A-02 tests)
@pytest.fixture()
def core_fixture_dir(
    tmp_path: Path,
    docx_known_content: Path,
    pptx_known_content: Path,
    pdf_special_chars: Path,
    pdf_with_boilerplate: Path,
    pdf_no_text_layer: Path,
    pdf_too_short: Path,
    unsupported_format_file: Path,
) -> Path:
    import shutil

    data_dir = tmp_path / "data_fixture"
    for folder in ("Обзоры", "Статьи", "Доклады"):
        (data_dir / folder).mkdir(parents=True, exist_ok=True)

    shutil.copy(docx_known_content, data_dir / "Обзоры" / docx_known_content.name)
    shutil.copy(pptx_known_content, data_dir / "Статьи" / pptx_known_content.name)
    shutil.copy(pdf_special_chars, data_dir / "Доклады" / pdf_special_chars.name)
    shutil.copy(pdf_with_boilerplate, data_dir / "Обзоры" / pdf_with_boilerplate.name)
    shutil.copy(pdf_no_text_layer, data_dir / "Статьи" / pdf_no_text_layer.name)
    shutil.copy(pdf_too_short, data_dir / "Доклады" / pdf_too_short.name)
    shutil.copy(unsupported_format_file, data_dir / "Обзоры" / unsupported_format_file.name)
    return data_dir
